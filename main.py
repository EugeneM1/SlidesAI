pip install openai python-pptx
import os
import openai
from pptx import Presentation
from pptx.util import Inches
import argparse

# Function to generate slide content using ChatGPT
def generate_slide_content(description):
    openai.api_key = os.getenv("OPENAI_API_KEY")
    
    prompt = (
        "Create a detailed outline for a PowerPoint presentation based on the following description:\n\n"
        f"{description}\n\n"
        "The outline should include slide titles and detailed bullet points for each slide."
    )
    
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=500,
        n=1,
        stop=None,
        temperature=0.7,
    )
    
    return response.choices[0].text.strip()

# Function to create a PowerPoint presentation
def create_presentation(slide_content):
    prs = Presentation()
    
    slides = slide_content.split('\n\n')
    for slide in slides:
        if ':' in slide:
            title, content = slide.split(':', 1)
            slide_layout = prs.slide_layouts[1]
            slide_obj = prs.slides.add_slide(slide_layout)
            title_obj = slide_obj.shapes.title
            content_obj = slide_obj.placeholders[1]

            title_obj.text = title.strip()
            content_obj.text = content.strip()
    
    return prs

def main(description, output_file):
    slide_content = generate_slide_content(description)
    prs = create_presentation(slide_content)
    prs.save(output_file)
    print(f"Slideshow saved as {output_file}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate a slideshow from a description using ChatGPT.")
    parser.add_argument('--description', type=str, required=True, help='Description of the slideshow content.')
    parser.add_argument('--output', type=str, default='output.pptx', help='Output PowerPoint file name.')
    
    args = parser.parse_args()
    main(args.description, args.output)
